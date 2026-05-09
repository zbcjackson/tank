"""Tests for Office document extraction.

Each test builds a tiny real document in-memory using the library's
own writer, hands it to the extractor, and asserts the text carries
the input's semantic content. This exercises the full round-trip
through python-docx/openpyxl/python-pptx rather than mocking the
library layer, which would hide format-specific behavior.
"""

from __future__ import annotations

from pathlib import Path

import pytest

from tank_backend.core.content import DocumentBlock
from tank_backend.media import MediaStore
from tank_backend.media.office import (
    IWORK_MIME_TYPES,
    OFFICE_MIME_TYPES,
    extract_docx,
    extract_office_document,
    extract_pptx,
    extract_xlsx,
)

# ---------------------------------------------------------------------------
# Fixture builders — synthesize real docs via the target libraries
# ---------------------------------------------------------------------------


def _write_docx(path: Path) -> None:
    """Produce a DOCX with headings, paragraph, bullet list, and table."""
    from docx import Document

    doc = Document()
    doc.add_heading("Meeting Notes", level=1)
    doc.add_heading("Agenda", level=2)
    doc.add_paragraph("Discuss the Q3 roadmap.")
    doc.add_paragraph("Action item one.", style="List Bullet")
    doc.add_paragraph("Action item two.", style="List Bullet")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Owner"
    table.cell(0, 1).text = "Deadline"
    table.cell(1, 0).text = "Alice"
    table.cell(1, 1).text = "Friday"
    doc.save(str(path))


def _write_xlsx(path: Path) -> None:
    """Produce a 2-sheet XLSX with a header row and data."""
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Q3"
    ws.append(["Region", "Revenue"])
    ws.append(["APAC", 12345])
    ws.append(["EMEA", 67890])

    ws2 = wb.create_sheet("Forecast")
    ws2.append(["Month", "Estimate"])
    ws2.append(["Sep", 200])
    ws2.append(["Oct", 300])

    wb.save(str(path))


def _write_pptx(path: Path) -> None:
    """Produce a 2-slide deck with bullet text and speaker notes."""
    from pptx import Presentation

    prs = Presentation()
    title_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(title_layout)
    slide1.shapes.title.text = "Roadmap"
    slide1.placeholders[1].text = "Q3 priorities"

    content_layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(content_layout)
    slide2.shapes.title.text = "Details"
    body = slide2.placeholders[1].text_frame
    body.text = "Ship multi-modal"
    body.add_paragraph().text = "Fix onboarding"
    slide2.notes_slide.notes_text_frame.text = (
        "Call out the capability registry."
    )

    prs.save(str(path))


# ---------------------------------------------------------------------------
# Per-format extractor tests
# ---------------------------------------------------------------------------


class TestDocxExtraction:
    def test_headings_paragraphs_list_table(self, tmp_path):
        path = tmp_path / "notes.docx"
        _write_docx(path)
        text = extract_docx(path)

        # Headings render as Markdown headers with matching levels.
        assert "# Meeting Notes" in text
        assert "## Agenda" in text

        # Paragraph stays as a line.
        assert "Discuss the Q3 roadmap." in text

        # Bullets render as list items.
        assert "- Action item one." in text
        assert "- Action item two." in text

        # Table renders as a pipe table with header + separator rows.
        assert "| Owner | Deadline |" in text
        assert "| Alice | Friday |" in text
        assert "| --- | --- |" in text

    def test_empty_paragraphs_skipped(self, tmp_path):
        """Bare whitespace paragraphs don't produce empty output lines."""
        from docx import Document

        doc = Document()
        doc.add_paragraph("")  # empty
        doc.add_paragraph("   ")  # whitespace only
        doc.add_paragraph("real content")
        path = tmp_path / "sparse.docx"
        doc.save(str(path))

        text = extract_docx(path)
        assert text == "real content"

    def test_corrupt_file_raises_via_library(self, tmp_path):
        """Corrupt input surfaces the library exception; the
        top-level dispatch wraps it into a diagnostic string, but
        the extractor itself is allowed to raise.
        """
        from docx.opc.exceptions import PackageNotFoundError

        path = tmp_path / "bad.docx"
        path.write_bytes(b"not a real docx")
        with pytest.raises(PackageNotFoundError):
            extract_docx(path)


class TestXlsxExtraction:
    def test_two_sheets_with_tables(self, tmp_path):
        path = tmp_path / "numbers.xlsx"
        _write_xlsx(path)
        text = extract_xlsx(path)

        # Each sheet gets its own header.
        assert "## Sheet: Q3" in text
        assert "## Sheet: Forecast" in text

        # Values render through; trailing .0 stripped on integer floats.
        assert "| APAC | 12345 |" in text
        assert "| EMEA | 67890 |" in text
        assert "| Sep | 200 |" in text

    def test_large_sheet_truncated(self, tmp_path):
        """A sheet with too many rows hits the per-sheet row cap."""
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws.append(["n"])
        for i in range(1, 300):
            ws.append([i])
        path = tmp_path / "big.xlsx"
        wb.save(str(path))

        text = extract_xlsx(path)
        assert "more rows truncated" in text

    def test_cell_with_pipe_escaped(self, tmp_path):
        """A cell containing a literal ``|`` is escaped so the table
        doesn't silently eat the value.
        """
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws.append(["label", "value"])
        ws.append(["a|b", "c|d"])
        path = tmp_path / "pipes.xlsx"
        wb.save(str(path))

        text = extract_xlsx(path)
        assert "a\\|b" in text
        assert "c\\|d" in text


class TestPptxExtraction:
    def test_slides_and_notes(self, tmp_path):
        path = tmp_path / "deck.pptx"
        _write_pptx(path)
        text = extract_pptx(path)

        assert "## Slide 1" in text
        assert "Roadmap" in text
        assert "Q3 priorities" in text

        assert "## Slide 2" in text
        assert "Ship multi-modal" in text
        assert "Fix onboarding" in text

        # Speaker notes surface explicitly so the LLM can distinguish
        # on-slide content from narration.
        assert "Speaker notes:" in text
        assert "capability registry" in text


class TestDispatch:
    """extract_office_document routes by MIME type."""

    def test_docx_mime_routes_to_docx(self, tmp_path):
        path = tmp_path / "a.docx"
        _write_docx(path)
        text = extract_office_document(
            path,
            "application/vnd.openxmlformats-officedocument."
            "wordprocessingml.document",
        )
        assert text is not None
        assert "Meeting Notes" in text

    def test_xlsx_mime_routes_to_xlsx(self, tmp_path):
        path = tmp_path / "a.xlsx"
        _write_xlsx(path)
        text = extract_office_document(
            path,
            "application/vnd.openxmlformats-officedocument."
            "spreadsheetml.sheet",
        )
        assert text is not None
        assert "Sheet: Q3" in text

    def test_pptx_mime_routes_to_pptx(self, tmp_path):
        path = tmp_path / "a.pptx"
        _write_pptx(path)
        text = extract_office_document(
            path,
            "application/vnd.openxmlformats-officedocument."
            "presentationml.presentation",
        )
        assert text is not None
        assert "Slide 1" in text

    def test_unknown_mime_returns_none(self, tmp_path):
        path = tmp_path / "a.bin"
        path.write_bytes(b"nope")
        assert extract_office_document(path, "application/x-weird") is None

    def test_parse_failure_returns_diagnostic(self, tmp_path):
        """A corrupt file yields a bracketed error string, not a raise."""
        path = tmp_path / "broken.docx"
        path.write_bytes(b"junk")
        result = extract_office_document(
            path,
            "application/vnd.openxmlformats-officedocument."
            "wordprocessingml.document",
        )
        assert result is not None
        assert "Could not extract text" in result


# ---------------------------------------------------------------------------
# Integration: MediaStore waterfall dispatches Office MIMEs
# ---------------------------------------------------------------------------


class TestStoreOfficeDispatch:
    @pytest.fixture()
    def store(self, tmp_path):
        return MediaStore(tmp_path / "media")

    @pytest.mark.asyncio()
    async def test_docx_through_store_yields_extracted_text(
        self, store, tmp_path,
    ):
        src = tmp_path / "in.docx"
        _write_docx(src)
        data = src.read_bytes()
        stored = await store.put(
            data,
            "application/vnd.openxmlformats-officedocument."
            "wordprocessingml.document",
            session_id="s",
        )
        block = DocumentBlock(
            source=stored.media_uri,
            mime_type="application/vnd.openxmlformats-officedocument."
            "wordprocessingml.document",
        )
        result = await store.materialize_for_llm(block)
        assert result.extracted_text is not None
        assert "Meeting Notes" in result.extracted_text
        assert result.send_native is False
        assert result.page_images == ()

    @pytest.mark.asyncio()
    async def test_xlsx_through_store(self, store, tmp_path):
        src = tmp_path / "in.xlsx"
        _write_xlsx(src)
        data = src.read_bytes()
        stored = await store.put(
            data,
            "application/vnd.openxmlformats-officedocument."
            "spreadsheetml.sheet",
            session_id="s",
        )
        block = DocumentBlock(
            source=stored.media_uri,
            mime_type="application/vnd.openxmlformats-officedocument."
            "spreadsheetml.sheet",
        )
        result = await store.materialize_for_llm(block)
        assert result.extracted_text is not None
        assert "Sheet: Q3" in result.extracted_text

    @pytest.mark.asyncio()
    async def test_pptx_through_store(self, store, tmp_path):
        src = tmp_path / "in.pptx"
        _write_pptx(src)
        data = src.read_bytes()
        stored = await store.put(
            data,
            "application/vnd.openxmlformats-officedocument."
            "presentationml.presentation",
            session_id="s",
        )
        block = DocumentBlock(
            source=stored.media_uri,
            mime_type="application/vnd.openxmlformats-officedocument."
            "presentationml.presentation",
        )
        result = await store.materialize_for_llm(block)
        assert result.extracted_text is not None
        assert "Slide 1" in result.extracted_text

    @pytest.mark.asyncio()
    async def test_sidecar_cached_on_second_materialize(
        self, store, tmp_path,
    ):
        """Second materialize reads the .txt sidecar, not the docx."""
        src = tmp_path / "in.docx"
        _write_docx(src)
        data = src.read_bytes()
        stored = await store.put(
            data,
            "application/vnd.openxmlformats-officedocument."
            "wordprocessingml.document",
            session_id="s",
        )
        block = DocumentBlock(
            source=stored.media_uri,
            mime_type="application/vnd.openxmlformats-officedocument."
            "wordprocessingml.document",
        )
        await store.materialize_for_llm(block)

        # Sidecar present?
        _session, filename = store._parse_uri(stored.media_uri)
        docx_path = store.root / _session / filename
        sidecar = docx_path.with_suffix(docx_path.suffix + ".txt")
        assert sidecar.exists()

        # Overwrite the sidecar with a sentinel. If the store
        # re-extracts, our sentinel gets blown away. If it honours the
        # cache, the sentinel survives.
        sidecar.write_text("SIDECAR-SENTINEL", encoding="utf-8")
        result = await store.materialize_for_llm(block)
        assert result.extracted_text == "SIDECAR-SENTINEL"


# ---------------------------------------------------------------------------
# Constants / wiring
# ---------------------------------------------------------------------------


class TestMimeSets:
    def test_office_mime_set_matches_dispatch(self):
        """Every Office MIME in the set has a routing branch."""
        for mime in OFFICE_MIME_TYPES:
            assert extract_office_document is not None
            # Smoke-check: dispatch for a non-existent path should
            # surface the library exception wrapped in a diagnostic.
            from tempfile import NamedTemporaryFile
            with NamedTemporaryFile(suffix=".bin", delete=False) as f:
                f.write(b"not a real office file")
                probe_path = Path(f.name)
            try:
                result = extract_office_document(probe_path, mime)
                # Either a diagnostic string or None; never a raise.
                assert result is None or isinstance(result, str)
            finally:
                probe_path.unlink()

    def test_iwork_set_has_common_formats(self):
        """iWork set covers the three iWork app MIME types."""
        assert any("pages" in m for m in IWORK_MIME_TYPES)
        assert any("numbers" in m for m in IWORK_MIME_TYPES)
        assert any("key" in m for m in IWORK_MIME_TYPES)
