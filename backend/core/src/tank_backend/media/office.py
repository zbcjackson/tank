"""Text extractors for Microsoft Office formats.

Each extractor converts an uploaded Office document into Markdown-ish
text the LLM can read. We go text-only for now; vision-model rendering
(DOCX → PDF → images) is deferred behind an explicit user request
because it needs a LibreOffice dependency we don't want to force on
every install.

Design choices worth knowing:

- **Markdown-ish output**, not pure prose. Headings become ``#``/``##``,
  tables become pipe tables, slide numbers are headers. This survives
  LLM context compression better than free-form paragraphs and gives
  the model cheap structural cues (it knows "Slide 3" refers to the
  third slide, not just "somewhere in the deck").

- **Output size bounded** by :data:`_MAX_CHARS`. A 500-sheet Excel can
  produce megabytes of text — eating the context window before the
  LLM has done anything. We truncate with a clear tail marker so the
  LLM knows there's more.

- **Lazy imports**. python-docx, openpyxl, python-pptx are all
  multi-MB. Importing them at module load would hurt cold-start for
  users who never upload Office files. We pay the cost only on the
  first upload of each format.
"""

from __future__ import annotations

import logging
from pathlib import Path

logger = logging.getLogger(__name__)


# MIME types we accept. The upload endpoint / WebSocket path decides
# whether to create a DocumentBlock by consulting this set.
OFFICE_MIME_TYPES: frozenset[str] = frozenset({
    # DOCX
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    # Legacy .doc deliberately omitted — python-docx only handles OOXML.
    # XLSX
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    # PPTX
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
})

# iWork MIME types we reject with a helpful error. Users sometimes end
# up sending these when AirDropping from a Mac to the web UI; the error
# message matters more than supporting them.
IWORK_MIME_TYPES: frozenset[str] = frozenset({
    "application/x-iwork-pages-sffpages",       # Pages
    "application/x-iwork-numbers-sffnumbers",   # Numbers
    "application/x-iwork-keynote-sffkey",       # Keynote
    # Pre-2013 variants — rare today but worth the two-line match.
    "application/vnd.apple.pages",
    "application/vnd.apple.numbers",
    "application/vnd.apple.keynote",
})

# Hard cap on extractor output. Bounds context spend; larger documents
# get truncated with an explanatory tail. Rounded to ~20K tokens
# assuming 4 chars/token — leaves plenty of room for the rest of the
# conversation plus the user's question about the document.
_MAX_CHARS = 80_000

# How many rows per sheet we include. Excel files with millions of
# rows are a reality; LLMs can't read that many anyway.
_MAX_ROWS_PER_SHEET = 200


def _truncate(text: str) -> str:
    """Apply :data:`_MAX_CHARS` with a marker. The marker tells the
    LLM there's more data so it can ask the user to narrow the query.
    """
    if len(text) <= _MAX_CHARS:
        return text
    return (
        text[:_MAX_CHARS]
        + f"\n\n[... truncated after {_MAX_CHARS} characters; upload "
        "a narrower selection to see the rest.]"
    )


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------


def extract_docx(path: Path) -> str:
    """Convert a .docx file to Markdown-ish text.

    Captures paragraphs (with heading levels when present), tables
    (as pipe tables), and bullet/numbered list items. Inline images
    and text boxes are skipped — python-docx can read them but they
    have no meaningful text representation at this scope.
    """
    from docx import Document
    from docx.oxml.ns import qn

    doc = Document(str(path))
    out: list[str] = []

    # Iterate body elements in document order so headings, paragraphs,
    # and tables interleave naturally. The document.paragraphs /
    # document.tables accessors silo them, which loses ordering.
    body = doc.element.body
    for child in body.iterchildren():
        tag = child.tag
        if tag == qn("w:p"):
            para = _paragraph_from_element(doc, child)
            if para is None:
                continue
            text = para.text.strip()
            if not text:
                continue
            style_obj = para.style
            style_name = style_obj.name if (style_obj and style_obj.name) else ""
            style = style_name.lower()
            if style.startswith("heading"):
                # "Heading 1" -> "# ", "Heading 2" -> "## ", …
                level = _heading_level(style)
                out.append(f"{'#' * level} {text}")
            elif style.startswith("list"):
                out.append(f"- {text}")
            else:
                out.append(text)
        elif tag == qn("w:tbl"):
            tbl = _table_from_element(doc, child)
            if tbl is not None:
                out.append(_render_docx_table(tbl))

    return _truncate("\n\n".join(out).strip())


def _paragraph_from_element(doc, element):
    """Find the Paragraph object whose underlying XML is ``element``."""
    # python-docx caches mapping internally; the cheapest safe way is
    # to walk document.paragraphs since the list order matches body
    # order. Fall through to None if the element came from a
    # non-standard container we don't want to traverse.
    for p in doc.paragraphs:
        if p._element is element:
            return p
    return None


def _table_from_element(doc, element):
    for t in doc.tables:
        if t._element is element:
            return t
    return None


def _heading_level(style_name: str) -> int:
    """Pick the heading level from a style name like "heading 3"."""
    parts = style_name.split()
    if len(parts) >= 2 and parts[-1].isdigit():
        level = int(parts[-1])
        return max(1, min(level, 6))
    return 1


def _render_docx_table(table) -> str:
    """Render a docx table as a Markdown pipe table.

    python-docx tables can have merged cells; we emit the raw text
    from each cell without trying to reconstruct the merge structure.
    That's honest — the LLM sees repeated content which is better
    than silently losing data.
    """
    rows: list[list[str]] = []
    for row in table.rows:
        cells = [cell.text.strip().replace("|", "\\|") or " " for cell in row.cells]
        rows.append(cells)
    if not rows:
        return ""
    header = rows[0]
    sep = ["---"] * len(header)
    body = rows[1:] if len(rows) > 1 else []
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(sep) + " |",
        *("| " + " | ".join(r) + " |" for r in body),
    ]
    return "\n".join(lines)


# ---------------------------------------------------------------------------
# XLSX
# ---------------------------------------------------------------------------


def extract_xlsx(path: Path) -> str:
    """Convert a .xlsx workbook into per-sheet Markdown tables.

    We load with ``data_only=True`` so cells show computed values, not
    formulas. The formula form is often useless to an LLM without the
    evaluation context. If we need formulas later, pass the parameter
    through from config.
    """
    from openpyxl import load_workbook

    wb = load_workbook(str(path), data_only=True, read_only=True)
    out: list[str] = []

    for sheet in wb.worksheets:
        out.append(f"## Sheet: {sheet.title}")
        rows_rendered = 0
        table_rows: list[list[str]] = []

        # read_only workbooks use a generator for rows; iterate once.
        for row_idx, row in enumerate(sheet.iter_rows(values_only=True), start=1):
            # Strip trailing empty cells so we don't emit a forest of
            # "|  |  |  |" columns from sparse data.
            cells = list(row)
            while cells and (cells[-1] is None or cells[-1] == ""):
                cells.pop()
            if not cells:
                continue

            formatted = [_format_cell(c) for c in cells]
            table_rows.append(formatted)
            rows_rendered += 1
            if rows_rendered >= _MAX_ROWS_PER_SHEET:
                out.append(_render_xlsx_table(table_rows))
                out.append(
                    f"[... {sheet.max_row - row_idx} more rows truncated.]"
                )
                break
        else:
            if table_rows:
                out.append(_render_xlsx_table(table_rows))

    wb.close()
    return _truncate("\n\n".join(out).strip())


def _format_cell(value) -> str:
    """Coerce a cell value to a short string, escaping pipe characters.

    Dates render via ``isoformat()``; numbers drop trailing ``.0`` so a
    cell entered as ``42`` doesn't become ``42.0`` on read-back.
    """
    if value is None:
        return " "
    if isinstance(value, float) and value.is_integer():
        value = int(value)
    text = str(value).strip()
    if not text:
        return " "
    return text.replace("|", "\\|").replace("\n", " ")


def _render_xlsx_table(rows: list[list[str]]) -> str:
    """Render rows as a Markdown pipe table.

    First row is treated as the header. Columns are normalised to the
    widest row so the table is rectangular — narrow rows get padded.
    """
    width = max(len(r) for r in rows)
    rows = [r + [" "] * (width - len(r)) for r in rows]
    header = rows[0]
    sep = ["---"] * width
    body = rows[1:] if len(rows) > 1 else []
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(sep) + " |",
        *("| " + " | ".join(r) + " |" for r in body),
    ]
    return "\n".join(lines)


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------


def extract_pptx(path: Path) -> str:
    """Convert a .pptx deck to slide-by-slide markdown.

    Each slide becomes a ``## Slide N`` section. We collect text from
    shapes in placement order and append any speaker notes at the end
    of the slide's section — notes often have the most substantive
    content (the actual narrative the presenter wrote).
    """
    from pptx import Presentation

    prs = Presentation(str(path))
    out: list[str] = []

    for i, slide in enumerate(prs.slides, 1):
        out.append(f"## Slide {i}")
        shape_texts: list[str] = []
        for shape in slide.shapes:
            # has_text_frame is defined on Shape (not BaseShape) — the
            # common superclass hides it, so narrow via getattr rather
            # than trusting the static type.
            if not getattr(shape, "has_text_frame", False):
                continue
            text_frame = getattr(shape, "text_frame", None)
            if text_frame is None:
                continue
            text = text_frame.text.strip()
            if text:
                shape_texts.append(text)
        if shape_texts:
            out.append("\n".join(shape_texts))

        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            notes = notes_frame.text.strip() if notes_frame else ""
            if notes:
                out.append(f"**Speaker notes:** {notes}")

    return _truncate("\n\n".join(out).strip())


# ---------------------------------------------------------------------------
# Dispatch
# ---------------------------------------------------------------------------


def extract_office_document(path: Path, mime_type: str) -> str | None:
    """Route to the right extractor by MIME type.

    Returns extracted text, or ``None`` if the MIME isn't handled.
    Never raises — a parsing failure returns a short diagnostic so the
    LLM sees something useful even on a corrupt upload.
    """
    try:
        if mime_type == (
            "application/vnd.openxmlformats-officedocument."
            "wordprocessingml.document"
        ):
            return extract_docx(path)
        if mime_type == (
            "application/vnd.openxmlformats-officedocument."
            "spreadsheetml.sheet"
        ):
            return extract_xlsx(path)
        if mime_type == (
            "application/vnd.openxmlformats-officedocument."
            "presentationml.presentation"
        ):
            return extract_pptx(path)
    except Exception as exc:  # pragma: no cover — defensive
        logger.warning("Office extraction failed (%s): %s", mime_type, exc)
        return f"[Could not extract text from {mime_type}: {exc}]"
    return None


__all__ = [
    "IWORK_MIME_TYPES",
    "OFFICE_MIME_TYPES",
    "extract_docx",
    "extract_office_document",
    "extract_pptx",
    "extract_xlsx",
]
